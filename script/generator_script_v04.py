import fitz  # PyMuPDF
import os
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import MSO_AUTO_SIZE

def extract_specific_images(pdf_path):
    print(f"Képek kinyerése a {pdf_path} fájlból...")
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"Nem található a fájl: {pdf_path}")

    doc = fitz.open(pdf_path)
    
    # Oldalak hozzárendelése a diákhoz (1-alapú indexelés)
    slide_to_page_map = {
        0: 5,   # Slide 1: Global production communities
        1: 4,   # Slide 2: Giant pile of EF lenses
        2: 10,  # Slide 3: EF vs RF mount comparison
        3: 12,  # Slide 4: Aberration ray bundle
        4: 14,  # Slide 5: Control Ring layout
        5: 16,  # Slide 6: RF28-70mm F2 lens
        6: 24,  # Slide 7: RF24-105mm F4 lens
        7: 34,  # Slide 8: EF-EOS R Drop-in Filter
        8: 35,  # Slide 9: 5-axis Image Stabilization
        9: 41   # Slide 10: Flange back distance comparison
    }
    
    slide_images = {}
    
    for slide_idx, page_num in slide_to_page_map.items():
        if page_num <= len(doc):
            page = doc[page_num - 1] # 0-indexed
            image_list = page.get_images(full=True)
            
            if image_list:
                # Kinyerjük az első képet az adott oldalon (általában ez a fő ábra)
                xref = image_list[0][0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                
                img_filename = f"slide_{slide_idx+1}_page_{page_num}.{image_ext}"
                with open(img_filename, "wb") as f:
                    f.write(image_bytes)
                
                slide_images[slide_idx] = img_filename
                print(f"Dia {slide_idx+1}: Kép kinyerve a(z) {page_num}. oldalról.")
            else:
                print(f"Dia {slide_idx+1}: Nem található kép a(z) {page_num}. oldalon.")
        else:
            print(f"Dia {slide_idx+1}: Az oldalszám ({page_num}) meghaladja a PDF hosszát.")
            
    return slide_images

def create_presentation():
    pdf_filename = "canon-eos-r-white-paper.pdf"
    
    # 1. Képek kinyerése a PDF-ből
    slide_images = extract_specific_images(pdf_filename)

    # 2. Prezentáció adatszerkezete
    slides_data = [
        {
            "title": "A Canon EOS R Rendszer: Új Korszak",
            "bullets": [
                "Célkitűzés: Az optikai teljesítmény és a kreatív működési rugalmasság maximalizálása a globális piacon.",
                "Fókuszterület: Új generációs, full-frame tükör nélküli (mirrorless) technológia bevezetése.",
                "A Rendszer Lényege: Teljesen új tervezésű, innovatív RF bajonett (lens mount) architektúra.",
                "Kompatibilitás: A meglévő EF és EF-S lencsék teljes körű támogatásának integrálása."
            ],
            "notes": "Üdvözlöm önöket! A mai prezentációban áttekintjük a Canon EOS R rendszerének stratégiai és műszaki alapjait."
        },
        {
            "title": "A Harmincéves EF Rendszer Korlátai",
            "bullets": [
                "Fizikai gátak: A hagyományos EF bajonett (54 mm átmérő, 44 mm távolság) elérte a rugalmasság határait.",
                "Adatátviteli sebesség: A 8-tűs (pin) elektronikus kapcsolat sávszélessége mára már korlátozó tényező.",
                "Fejlesztési limitációk: Kevés az elektronikus csatorna az összetett operációs funkciókhoz.",
                "Fókusz korlátok: Szűkös lehetőségek a modern, szenzor-alapú autofókusz (AF) rendszerekben."
            ],
            "notes": "Bár a három évtizedes EF rendszer hatalmas siker volt, az új jövőbeni igények új architektúrát követeltek."
        },
        {
            "title": "Az RF Bajonett Innovációja",
            "bullets": [
                "Belső átmérő megőrzése: Az 54 mm-es átmérő megtartása az optimális fényáteresztés érdekében.",
                "Rövidített bázistávolság: A szenzor és a bajonett távolsága 44 mm-ről mindössze 20 mm-re csökkent.",
                "Növelt sávszélesség: A csatlakozók számának 8-ról 12-re történő emelése biztosítja a nagysebességű kommunikációt.",
                "Tervezési szabadság: Nagy átmérőjű hátsó lencsetagok helyezhetők közvetlenül a full-frame szenzor elé."
            ],
            "notes": "Ez a dia a rendszer szíve. A 12 aranyozott érintkező és a rövid 20 mm-es távolság sosem látott szabadságot ad."
        },
        {
            "title": "Optikai Aberrációk Kezelése",
            "bullets": [
                "Fénytörési kihívások: A fénytörés fokozásával drasztikusan nőnek a kromatikus és monokromatikus képhibák.",
                "Optikai megoldás: A nagy hátsó lencsék 'kíméletesebb' szögben vetítik a fénysugarakat a szenzor sarkaiba.",
                "Digital Lens Optimizer (DLO): Beépített képoptimalizáló szoftverrendszer a kamerában a finomfelbontás javítására.",
                "DLO Működése: Valós időben, a RAW fájlok feldolgozásakor korrigálja a diffrakciót és a képhibákat."
            ],
            "notes": "Az új RF rendszer lencsetagjai csökkentik a peremtorzítást, a DLO pedig beépített szoftveres motorral tökéletesíti a képet."
        },
        {
            "title": "Az RF Lencsék Exkluzív Funkciói",
            "bullets": [
                "Vezérlőgyűrű (Control Ring): A fókuszgyűrű mellett egy új, testreszabható gyűrű az expozíciós beállításokhoz.",
                "Változtatható fókusz-irány: A kézi élességállítás gyűrűjének forgásiránya szoftveresen megfordítható.",
                "'Fly-by-wire' fókuszálás: Nincs mechanikus kapcsolat; a fókuszgyűrű csak elektronikus jeleket küld a motornak.",
                "Szenzorvédelem: Kikapcsoláskor a rekeszlamellák automatikusan bezárulnak, védve a szenzort."
            ],
            "notes": "A Control Ring ergonómiai áttörés, a fly-by-wire fókusz pedig zökkenőmentes videós átmeneteket garantál."
        },
        {
            "title": "Az Induló RF Lencsekínálat (1. Rész)",
            "bullets": [
                "RF28-70mm F2 L USM: Szabványos, L-szériás zoomobjektív egyedülálló, állandó f/2.0 fényerővel.",
                "Méretbeli előny: A nagy bajonett és rövid bázistávolság miatt a frontlencse része kompaktabb maradhatott.",
                "RF50mm F1.2 L USM: Nagy átmérőjű, standard fix objektív extrém f/1.2 rekeszértékkel.",
                "Fejlett bevonatok: Air Sphere Coating technológia a becsillanások és szellemkép minimálisra csökkentésére."
            ],
            "notes": "A 28-70-es f/2-es zoom megalkotása az EF rendszerben fizikai képtelenség lett volna ilyen paraméterekkel."
        },
        {
            "title": "Az Induló RF Lencsekínálat (2. Rész)",
            "bullets": [
                "RF24-105mm F4 L IS USM: Kompakt f/4-es zoom, integrált képstabilizátorral professzionális feladatokhoz.",
                "Nano USM meghajtás: Extra vékony, rendkívül gyors és néma ultrahangos fókuszmotor, ideális videózáshoz.",
                "RF35mm F1.8 MACRO IS STM: Széles látószögű objektív f/1.8 fényerővel és 0.5x makró nagyítással.",
                "Kompakt felépítés: Léptetőmotor (STM) gondoskodik a finom mozgásról, könnyű kivitelben."
            ],
            "notes": "Ez a két objektív a rendszer 'mindenese', kiváló és hangtalan fókuszálással."
        },
        {
            "title": "Rendszerkompatibilitás: EF Mount Adapterek",
            "bullets": [
                "Stratégiai integráció: A hatalmas globális EF/EF-S lencsekínálat maradéktalan átemelése az új rendszerbe.",
                "Mount Adapter EF-EOS R: Alapmodell a tökéletes fizikai és elektronikus csatlakoztatáshoz.",
                "Control Ring Adapter: A régebbi EF lencséket is felruházza az új, beépített expozíció-vezérlő gyűrűvel.",
                "Drop-In Filter Adapter: Professzionális adapter beépített, cserélhető szűrő-nyílással (cirkuláris polár vagy ND)."
            ],
            "notes": "Az adapterek révén a több évtizedes EF objektívek még több funkciót kapnak, mint eredeti vázukon."
        },
        {
            "title": "Képstabilizáció és Dual Pixel AF",
            "bullets": [
                "Interaktív Stabilizáció: A lencse giroszkópja és a szenzor adatai közösen harcolnak a remegés ellen.",
                "5-Tengelyes videós IS: A lencse optikai és a szenzor elektronikus stabilizációja öttengelyessé egészül ki.",
                "Dual Pixel fókuszálás: Minden egyes képalkotó szenzorpixelben két fotodióda végzi a fázisérzékelést.",
                "Kiterjesztett lefedettség: A választható fókuszpontok a kép szélességének 88%-át, magasságának 100%-át lefedik."
            ],
            "notes": "A DIGIC 8 processzor valós időben dolgozza fel a lencse és a szenzor adatait az 5-tengelyes kompenzációhoz."
        },
        {
            "title": "Összegzés: A Jövő Alapköve",
            "bullets": [
                "Mérnöki paradigmaváltás: Az RF rendszer az EF hagyatékára építve bontja le a lencsetervezés fizikai korlátait.",
                "Szinergia: Tökéletesebb hardveres integráció az optika, a mechanika és a digitális jelfeldolgozás között.",
                "Kész a jövő kihívásaira: Potenciál a folyamatosan növekvő szenzorfelbontáshoz és fényerő-igényekhez.",
                "Kompakt forma: Jelentős méretcsökkenés anélkül, hogy a strukturális stabilitás vagy a teljesítmény csorbulna."
            ],
            "notes": "A Canon EOS R rendszer évtizedekre biztosítja az alapot a legmagasabb szintű fotós és videós munkákhoz. Köszönöm a figyelmet!"
        }
    ]

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for i, slide_data in enumerate(slides_data):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        # --- 1. CÍM ---
        title_shape = slide.shapes.title
        title_shape.text = "" 
        
        # Szigorúan fix pozíció, hogy esélye se legyen rálógni a szövegre
        title_shape.left = Inches(0.5)
        title_shape.top = Inches(0.4)
        title_shape.width = Inches(9.0)
        title_shape.height = Inches(0.8)
        
        p_title = title_shape.text_frame.paragraphs[0]
        run_title = p_title.add_run()
        run_title.text = slide_data["title"]
        run_title.font.name = 'Arial'
        run_title.font.size = Pt(32)
        run_title.font.bold = True

        # --- 2. VÁZLAT PONTOK ---
        body_shape = slide.placeholders[1]
        
        # Teljesen elkülönített pozíció a címtől (top = 1.6 inch) és a képtől (width = 4.8 inch)
        body_shape.left = Inches(0.5)
        body_shape.top = Inches(1.6)
        body_shape.width = Inches(4.8)
        body_shape.height = Inches(5.2)
        
        tf = body_shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        for bullet in slide_data["bullets"]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.name = 'Calibri'
            p.font.size = Pt(18)
            p.space_after = Pt(14) 

        # --- 3. KONTEXTUS-AZONOS KÉP ---
        if i in slide_images:
            img_path = slide_images[i]
            # Kép pozíciója a jobb oldalon: balról 5.5 inch (így van 0.2 inch rés a szöveg és kép között)
            # A képarány megőrzéséhez csak a szélességet (width) rögzítjük
            try:
                slide.shapes.add_picture(img_path, left=Inches(5.5), top=Inches(1.6), width=Inches(4.0))
            except Exception as e:
                print(f"Hiba a(z) {img_path} beillesztésekor a(z) {i+1}. dián: {e}")

        # --- 4. ELŐADÓI JEGYZET ---
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = slide_data["notes"]

    filename = "Canon_EOS_R_Executive_Summary_v04.pptx"
    prs.save(filename)
    print(f"\nA prezentáció sikeresen mentve lett: {filename}")

if __name__ == "__main__":
    create_presentation()
