import os
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE

# A 10 dia tartalma és a hivatkozott PDF oldal indexek (0-indexált)
slides_data = [
    {
        "title": "1. Bevezetés a Canon EOS R Rendszerbe",
        "bullets": [
            "A három évtizedes EOS történelem legújabb technológiai mérföldköve[cite: 1208, 1221].",
            "Válasz az álló- és mozgóképkészítés határainak elmosódására[cite: 1210].",
            "A tükör nélküli (mirrorless) technológia növekvő dominanciája[cite: 1214].",
            "Kibővített optikai tervezési szabadság a jövő fejlesztéseihez[cite: 1215]."
        ],
        "notes": "Tisztelt Vezetőség! A mai prezentációban a Canon EOS R rendszer mérnöki és üzleti jelentőségét vizsgáljuk meg. Harminc év sikeres EF rendszer után a piac paradigmaváltást követel. Az EOS R nem csupán egy új kamera, hanem egy teljesen új, jövőtálló optikai ökoszisztéma[cite: 1208, 1221].",
        "page_idx": 3
    },
    {
        "title": "2. A Jelenlegi EF Rendszer Korlátai",
        "bullets": [
            "Az 54 mm-es belső átmérő kiváló, de a 44 mm-es bázistávolság korlátozó tényező[cite: 1325, 1326].",
            "Elégtelen adatátviteli sebesség a modern szenzorok igényeihez[cite: 1257].",
            "Kommunikációs csatornák szűkössége az új operatív funkciókhoz[cite: 1258].",
            "Korlátok a szenzoralapú autofókusz (AF) rendszerek maximalizálásában[cite: 1259]."
        ],
        "notes": "Bár az EF rendszer hosszú ideig iparági standard volt, a 44 milliméteres szenzortávolság fizikai határt szabott az aszférikus és nagy rekeszű lencsék integrálásának[cite: 1325, 1326]. Továbbá a meglévő nyolc érintkezős elektronikai sávszélesség nem képes kiszolgálni a legújabb hibrid fókuszrendszereket[cite: 1257].",
        "page_idx": 4
    },
    {
        "title": "3. Az Ideális Optikai Rendszer Tervezési Modellje",
        "bullets": [
            "Folyamatosan növekvő full-frame szenzor felbontás adaptációja[cite: 1273, 1274].",
            "A dinamikatartomány és az expozíciós határok kitolása[cite: 1275].",
            "A méret/súly, az optikai érzékenység és az operatív képességek optimális egyensúlya[cite: 1281].",
            "Több paraméteres optimalizáció a kreatív kompromisszumok elkerüléséért[cite: 1299]."
        ],
        "notes": "A lencsetervezés egy klasszikus 'trade-off' háromszögön alapul: méret, optikai minőség és funkcionalitás. Az új rendszer célja, hogy ennek a háromszögnek a területét megnövelje, ezáltal lehetővé téve kompromisszummentes megoldásokat, például magas fényerejű, de mégis hordozható zoom objektíveket[cite: 1271, 1281, 1299].",
        "page_idx": 7
    },
    {
        "title": "4. Az Új RF Bajonett Architektúra",
        "bullets": [
            "Megtartott 54 mm-es átmérő, mindössze 20 mm-es bázistávolsággal[cite: 1342, 1344].",
            "12 pines elektronikus csatlakozó a robusztus adatátvitelért[cite: 1345, 1460].",
            "Nagy átmérőjű hátsó lencsetagok elhelyezésének lehetősége[cite: 1370].",
            "Optimalizált 3-füles bajonett rögzítés, megújult pozicionálással[cite: 1364, 1365]."
        ],
        "notes": "A fizikai architektúra kulcsa a 20 mm-es Flange Back távolság. A mérnöki bravúr abban rejlik, hogy a hatalmas hátsó lencsetagok extrém közel kerülhetnek a szenzorhoz, ami szignifikánsan növeli a peremsötétedés és az aberrációk feletti kontrollt[cite: 1370]. A 12 pines csatoló jövőbiztos sávszélességet ad[cite: 1345].",
        "page_idx": 10
    },
    {
        "title": "5. Optikai Aberrációk Fizikai Kezelése",
        "bullets": [
            "A rövid bázistávolság enyhébb fénytörési szöget tesz lehetővé a peremeken[cite: 1412].",
            "Seidel-féle monokromatikus aberrációk (pl. kóma, asztigmatizmus) drasztikus csökkentése[cite: 1390, 1391].",
            "Diszperziós és kromatikus eltérések szoftveres-hardveres korrekciója[cite: 1394, 1396].",
            "Kiváló képélesség (MTF) a szenzor sarkaiban is[cite: 1400]."
        ],
        "notes": "A hagyományos lencséknél a peremre érkező fénysugarak meredek beesési szöge komoly torzításokat (aberrációkat) okoz[cite: 1398, 1399]. A 20 mm-es bázistávolságnak köszönhetően a fénysugarak sokkal tompább szögben, kíméletesebben érkeznek az érzékelő síkjára, jelentősen növelve az MTF értékeket a kép sarkaiban[cite: 1412].",
        "page_idx": 12
    },
    {
        "title": "6. Forradalmi RF Objektív Koncepciók",
        "bullets": [
            "Dedikált, programozható 'Control Ring' az expozíciós paraméterekhez[cite: 1430, 1431].",
            "Fókuszgyűrű forgásirányának szoftveres megfordíthatósága[cite: 1450, 1456].",
            "Mechanikus helyett nagyfelbontású elektronikus fókuszvezérlés (Focus-by-wire)[cite: 1451, 1453].",
            "Automatikus rekesz-zárás áramtalanításkor a szenzor védelmében[cite: 1445]."
        ],
        "notes": "Az RF objektívek mikrokontrollerei önálló egységként működnek. Az elektronikus 'Focus-by-wire' rendszer és a dedikált Control Ring soha nem látott ergonómiai flexibilitást biztosít az operatőrök és fotósok számára[cite: 1430, 1451]. Külön kiemelendő a kikapcsoláskor fellépő rekesz-záródás, amely megvédi a nyitott szenzort a lézerektől vagy a közvetlen naptól[cite: 1445].",
        "page_idx": 14
    },
    {
        "title": "7. Kiemelt RF Objektívek Teljesítménye",
        "bullets": [
            "RF28-70mm F2 L USM: Példátlan fényerő egy teljes zoomtartományon[cite: 1478, 1508].",
            "RF50mm F1.2 L USM: Kiemelkedő élesség aszférikus (ASC) bevonatokkal[cite: 1716, 1744].",
            "RF24-105mm F4 L IS USM: Kompakt méret, csendes Nano USM motor videókhoz[cite: 1833, 1864].",
            "RF35mm F1.8 MACRO: Kompakt széleslátószög, dedikált makró képességekkel[cite: 2043, 2046]."
        ],
        "notes": "A bemutatkozó lencsepark a bajonett határait demonstrálja. A 28-70-es F2-es zoom optikailag korábban lehetetlen volt ekkora méretben[cite: 1508]. A 24-105-ös modell az új vékonyított Nano USM motorjával pedig a videós produkciók professzionális igáslova[cite: 1864].",
        "page_idx": 16
    },
    {
        "title": "8. EF Visszafelé Kompatibilitás és Adapterek",
        "bullets": [
            "Veszteségmentes kommunikáció a meglévő EF/EF-S objektívekkel[cite: 2153, 2156].",
            "Standard EF-EOS R adapter a közvetlen mechanikai/elektronikai csatolásért[cite: 2156].",
            "Vezérlőgyűrűs (Control Ring) adapter a funkciók EF lencsékre való kiterjesztéséhez[cite: 2158].",
            "Belső szűrős (Drop-In Filter) adapter professzionális videósoknak (pl. V-ND)[cite: 2161, 2162]."
        ],
        "notes": "A beruházásvédelem kritikus szempont a B2B piacon. A Canon nem cserbenhagyja, hanem feljavítja a meglévő EF objektívparkot[cite: 2153]. A Control Ring és a beépített Drop-In szűrős adapterek olyan funkcionalitást adnak a régi lencséknek, amikkel korábban nem rendelkeztek[cite: 2158, 2161].",
        "page_idx": 31
    },
    {
        "title": "9. Dual Pixel CMOS AF és Képstabilizálás",
        "bullets": [
            "Fázisérzékelés dedikáltan, közvetlenül a szenzor fotodiódáin[cite: 2297, 2298].",
            "DIGIC 8 processzor a valós idejű, masszív adatfeldolgozáshoz[cite: 2210].",
            "Az objektív giroszkópja és a szenzor elmozdulás adatai szinkronizálnak[cite: 2209].",
            "Kombinált (optikai és elektronikus) 5-tengelyes stabilizálás[cite: 2226]."
        ],
        "notes": "A hardveres szinergia itt éri el a csúcsot. A lencsébe épített dual giroszkóp adatai összeolvadnak a szenzor Dual Pixel fázis-eltolódásos méréseivel a DIGIC 8 processzorban[cite: 2209, 2210]. Ez az adatfúzió 5-tengelyes stabilizálást és pengeéles videós követőfókuszt tesz lehetővé[cite: 2226].",
        "page_idx": 34
    },
    {
        "title": "10. Digital Lens Optimizer (DLO) Szoftveres Motor",
        "bullets": [
            "Hardveres szinten futó valós idejű lencsekorrekciós algoritmusok[cite: 2240].",
            "Diffrakció és kumulatív aberrációk azonnali kompenzálása[cite: 2278, 2279].",
            "A lencse teljes rekesztartománya kompromisszum nélkül használható[cite: 2289, 2290].",
            "Kiváló JPEG kimenet egyenesen a kamerából, csökkentett utómunka igény[cite: 2293]."
        ],
        "notes": "A Digital Lens Optimizer a modern jelfeldolgozás csúcsa. Ismeri az adott lencse Point Spread Function (PSF) adatait, és hardveres szinten, késleltetés nélkül vonja ki a képletből az optikai tökéletlenségeket[cite: 2240, 2250]. Ezzel a mérnökeink elérték, hogy az alkotóknak ne kelljen tartaniuk a diffrakciótól lerekeszelt állapotban[cite: 2289].",
        "page_idx": 37
    }
]

def create_executive_presentation(pdf_path, output_pptx):
    # PPTX inicializálása
    prs = Presentation()

    # PDF megnyitása a képek kinyeréséhez
    try:
        doc = fitz.open(pdf_path)
    except Exception as e:
        print(f"Hiba a PDF megnyitásakor: {e}")
        return

    for slide_info in slides_data:
        # Üres dia hozzáadása (index 6 a beépített sablonoknál)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Cím (Fix pozíció: Top 0.4", Left 0.5")
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9.0), Inches(0.8))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = slide_info["title"]
        title_p.font.name = 'Arial'
        title_p.font.bold = True
        title_p.font.size = Pt(28)

        # Vázlatpontok (Fix pozíció: Top 1.6", Left 0.5", Width 4.8")
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.8), Inches(5.0))
        body_frame = body_box.text_frame
        body_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        body_frame.word_wrap = True

        for point in slide_info["bullets"]:
            p = body_frame.add_paragraph()
            p.text = f"• {point}"
            p.level = 0
            p.font.name = 'Calibri'
            p.font.size = Pt(18)
            p.space_after = Pt(14)

        # Előadói jegyzet hozzáadása
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = slide_info["notes"]

        # Kép kinyerése és beillesztése (Jobb oldal: Top 1.6", Left 5.5")
        page_idx = slide_info["page_idx"]
        if page_idx < len(doc):
            page = doc[page_idx]
            image_list = page.get_images(full=True)
            
            if image_list:
                # A legnagyobb felbontású kép kiválasztása az apró grafikai elemek elkerülése végett
                xref = max(image_list, key=lambda img: img[2] * img[3])[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                img_ext = base_image["ext"]
                img_filename = f"temp_img_{page_idx}.{img_ext}"

                # Ideiglenes mentés
                with open(img_filename, "wb") as f:
                    f.write(image_bytes)

                # Kép hozzáadása a PPTX-hez, arányos átméretezéssel (width 4.0")
                slide.shapes.add_picture(img_filename, Inches(5.5), Inches(1.6), width=Inches(4.0))

                # Takarítás
                os.remove(img_filename)

    # Prezentáció mentése
    prs.save(output_pptx)
    print(f"Kész! A prezentáció sikeresen elmentve: {output_pptx}")

if __name__ == "__main__":
    # A fájlneveket a futtatási környezethez igazítsd
    INPUT_PDF = "canon-eos-r-white-paper.pdf"
    OUTPUT_PPTX = "final_prompt_generator_script_case_1.pptx"
    
    if os.path.exists(INPUT_PDF):
        create_executive_presentation(INPUT_PDF, OUTPUT_PPTX)
    else:
        print(f"A forrás PDF nem található: {INPUT_PDF}")