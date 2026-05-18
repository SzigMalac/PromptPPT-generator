import fitz  # PyMuPDF
import os
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import MSO_AUTO_SIZE

def extract_images_from_pdf(pdf_path, num_images=10):
    """Kinyeri a megadott számú képet a PDF-ből és elmenti őket helyileg."""
    print(f"Képek kinyerése a {pdf_path} fájlból...")
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"Nem található a fájl: {pdf_path}")

    doc = fitz.open(pdf_path)
    extracted_images = []
    image_count = 0
    
    for page_index in range(len(doc)):
        page = doc[page_index]
        image_list = page.get_images(full=True)
        
        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            
            image_filename = f"slide_img_{image_count + 1}.{image_ext}"
            with open(image_filename, "wb") as f:
                f.write(image_bytes)
            
            extracted_images.append(image_filename)
            image_count += 1
            
            if image_count >= num_images:
                print(f"{num_images} kép sikeresen kimentve.")
                return extracted_images
                
    print(f"Figyelem: Csak {image_count} képet találtam a PDF-ben.")
    return extracted_images

def create_presentation():
    # 1. Képek kinyerése a PDF-ből
    pdf_filename = "canon-eos-r-white-paper.pdf"
    extracted_images = extract_images_from_pdf(pdf_filename, num_images=10)

    # 2. Prezentáció adatszerkezete
    slides_data = [
        {
            "title": "A Canon EOS R Rendszer: Új Korszak a Képalkotásban",
            "bullets": [
                "Célkitűzés: Az optikai teljesítmény és a kreatív működési rugalmasság maximalizálása a globális piacon.",
                "Fókuszterület: Új generációs, full-frame tükör nélküli technológia bevezetése.",
                "A Rendszer Lényege: Teljesen új tervezésű, innovatív RF bajonett architektúra.",
                "Kompatibilitás: A meglévő EF és EF-S lencsék teljes körű támogatásának integrálása."
            ],
            "notes": "Üdvözlöm önöket! A mai prezentációban áttekintjük a Canon EOS R rendszerének stratégiai és műszaki alapjait."
        },
        {
            "title": "A Harmincéves EF Rendszer Korlátai",
            "bullets": [
                "Fizikai gátak: A hagyományos EF bajonett (54 mm átmérő, 44 mm távolság) elérte határait.",
                "Adatátviteli sebesség: A 8-tűs elektronikus kapcsolat sávszélessége korlátozó tényező.",
                "Fejlesztési limitációk: Kevés elektronikus csatorna az összetett operációs funkciókhoz.",
                "Fókusz korlátok: Szűkös lehetőségek a modern szenzor-alapú autofókusz kiszolgálásában."
            ],
            "notes": "Bár az EF rendszer hatalmas siker volt, a jövőbeni igények új architektúrát követeltek."
        },
        {
            "title": "Az RF Bajonett Innovációja: A Fizika Újraírása",
            "bullets": [
                "Belső átmérő megőrzése: Az 54 mm-es átmérő megtartása az optimális fényáteresztésért.",
                "Rövidített bázistávolság: A távolság 44 mm-ről mindössze 20 mm-re csökkent.",
                "Növelt sávszélesség: 8 helyett 12 csatlakozó biztosítja a nagysebességű kommunikációt.",
                "Tervezési szabadság: Nagy hátsó lencsetagok helyezhetők közvetlenül a szenzor elé."
            ],
            "notes": "A 12 aranyozott érintkező hatalmas adatsűrűséget biztosít, a rövid bázistávolság pedig új optikai tervezési szabadságot ad."
        },
        {
            "title": "Optikai Aberrációk Kezelése",
            "bullets": [
                "Fénytörési kihívások: A fénytörés növelésével nőnek a kromatikus és monokromatikus képhibák.",
                "Optikai megoldás: A nagy hátsó lencsék kíméletesebb szögben vetítik a fényt a szélekre.",
                "Digital Lens Optimizer (DLO): Beépített képoptimalizáló szoftverrendszer a kamerában.",
                "DLO Működése: RAW fájloknál valós időben korrigálja a diffrakciót és a képhibákat."
            ],
            "notes": "Az új RF rendszer lencsetagjai csökkentik a peremtorzítást, a DLO pedig beépített szoftveres motorral tökéletesíti a képet."
        },
        {
            "title": "Új Felhasználói Élmény az RF Lencséken",
            "bullets": [
                "Vezérlőgyűrű (Control Ring): Új, testreszabható gyűrű expozíciós beállításokhoz.",
                "Változtatható fókusz-irány: A kézi élességállítás forgásiránya szoftveresen megfordítható.",
                "Fly-by-wire fókuszálás: Nincs mechanikus kapcsolat, elektronikus jelek vezérlik a motort.",
                "Szenzorvédelem: Kikapcsoláskor a rekeszlamellák védelmi célból automatikusan bezárulnak."
            ],
            "notes": "A Control Ring ergonómiai áttörés, a fly-by-wire fókusz pedig zökkenőmentes videós átmeneteket garantál."
        },
        {
            "title": "Az Induló RF Lencsekínálat (1. Rész)",
            "bullets": [
                "RF28-70mm F2 L USM: Szabványos zoomobjektív egyedülálló, állandó f/2.0 fényerővel.",
                "Méretbeli előny: Az új bajonett miatt a frontlencse része kompaktabb maradhatott.",
                "RF50mm F1.2 L USM: Nagy átmérőjű fix objektív extrém f/1.2 rekeszértékkel.",
                "Fejlett bevonatok: Air Sphere Coating technológia a becsillanások csökkentésére."
            ],
            "notes": "A 28-70-es f/2-es zoom megalkotása az EF rendszerben fizikai képtelenség lett volna ekkora méretben."
        },
        {
            "title": "Az Induló RF Lencsekínálat (2. Rész)",
            "bullets": [
                "RF24-105mm F4 L IS USM: Kompakt f/4-es zoom integrált képstabilizátorral.",
                "Nano USM meghajtás: Vékony, gyors és néma ultrahangos fókuszmotor videózáshoz.",
                "RF35mm F1.8 MACRO IS STM: Széles látószögű objektív 0.5x makró nagyítással.",
                "Kompakt felépítés: Léptetőmotor (STM) gondoskodik a finom mozgásról."
            ],
            "notes": "Ezek a lencsék rendkívül sokoldalúak, a Nano USM és az STM motorok pedig hangtalan és precíz fókuszálást tesznek lehetővé."
        },
        {
            "title": "Rendszerkompatibilitás: EF Mount Adapterek",
            "bullets": [
                "Stratégiai integráció: A hatalmas EF/EF-S lencsekínálat zökkenőmentes átemelése.",
                "Mount Adapter EF-EOS R: Alapmodell tökéletes fizikai/elektronikus csatlakoztatáshoz.",
                "Control Ring Adapter: A régebbi EF lencséket is felruházza az expozíció-vezérlő gyűrűvel.",
                "Drop-In Filter Adapter: Adapter beépített, cserélhető szűrő-nyílással (ND vagy Polár)."
            ],
            "notes": "Az adapterek révén a több évtizedes EF objektívek még több funkciót kapnak, mint eredeti vázukon."
        },
        {
            "title": "Képstabilizáció és Dual Pixel CMOS AF",
            "bullets": [
                "Interaktív Stabilizáció: Lencse giroszkóp és szenzor adatok közösen harcolnak a remegés ellen.",
                "5-Tengelyes videós IS: Optikai és elektronikus stabilizáció öttengelyes kombinációja.",
                "Dual Pixel fókusz: Minden szenzorpixelben két fotodióda végzi a fázisérzékelést.",
                "Kiterjesztett lefedettség: A választható fókuszpontok a kép 88x100%-át lefedik."
            ],
            "notes": "A DIGIC 8 processzor valós időben dolgozza fel a lencse és a szenzor adatait az 5-tengelyes kompenzációhoz."
        },
        {
            "title": "Összegzés: A Képalkotás Jövőjének Alapköve",
            "bullets": [
                "Mérnöki paradigmaváltás: Az RF rendszer lebontja az EF rendszer fizikai korlátait.",
                "Szinergia: Tökéletes integráció optika, mechanika és digitális jelfeldolgozás között.",
                "Kész a jövő kihívásaira: Potenciál növekvő felbontáshoz és extrém fényerő-igényekhez.",
                "Kompakt forma: Jelentős méretcsökkenés a strukturális stabilitás csorbulása nélkül."
            ],
            "notes": "A Canon EOS R rendszer évtizedekre biztosítja az alapot a legmagasabb szintű fotós és videós munkákhoz."
        }
    ]

    # 3. PPTX Létrehozása és formázása
    prs = Presentation()

    for i, slide_data in enumerate(slides_data):
        slide_layout = prs.slide_layouts[1] # Cím és tartalom layout
        slide = prs.slides.add_slide(slide_layout)

        # --- CÍM BEÁLLÍTÁSA (Arial) ---
        title_shape = slide.shapes.title
        title_shape.text = slide_data["title"]
        
        # Betűtípus Arial-ra cserélése a címben
        if title_shape.text_frame.paragraphs:
            for run in title_shape.text_frame.paragraphs[0].runs:
                run.font.name = 'Arial'

        # --- SZÖVEGDOBOZ POZICIONÁLÁSA ÉS FORMÁZÁSA (Calibri) ---
        body_shape = slide.placeholders[1]
        
        # Összeszűkítjük a bal oldali szövegdobozt, hogy elférjen a kép (szélesség ~ 4.8 inch)
        body_shape.width = Inches(4.8)
        
        tf = body_shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        for bullet in slide_data["bullets"]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            
            # Formázás: Calibri betűtípus és Fix 20 pt méret
            p.font.name = 'Calibri'
            p.font.size = Pt(20)

        # --- KÉP BEILLESZTÉSE A JOBB OLDALRA ---
        if i < len(extracted_images):
            img_path = extracted_images[i]
            # Kép pozicionálása: Balról 5.2 inchre (szövegdoboz után), fentről 1.8 inchre
            try:
                slide.shapes.add_picture(img_path, left=Inches(5.2), top=Inches(1.8), width=Inches(4.2))
            except Exception as e:
                print(f"Hiba a(z) {img_path} beillesztésekor: {e}")

        # --- ELŐADÓI JEGYZET BEÁLLÍTÁSA ---
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = slide_data["notes"]

    # 4. Fájl mentése
    filename = "Canon_EOS_R_Executive_Summary_Images.pptx"
    prs.save(filename)
    print(f"\nA prezentáció sikeresen elkészült és mentve lett: {filename}")
    
    # Opcionális: Kinyert képek törlése (ha nem akarod megtartani őket a mappában, vedd ki a kommentet)
    # for img in extracted_images:
    #     if os.path.exists(img):
    #         os.remove(img)

if __name__ == "__main__":
    create_presentation()