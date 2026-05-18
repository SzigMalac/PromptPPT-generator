from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE

def create_presentation():
    slides_data = [
        {
            "title": "A Canon EOS R Rendszer: Új Korszak a Képalkotásban",
            "bullets": [
                "Célkitűzés: Az optikai teljesítmény és a kreatív működési rugalmasság maximalizálása a globális piacon.",
                "Fókuszterület: Új generációs, full-frame tükör nélküli (mirrorless) technológia bevezetése álló- és mozgóképhez.",
                "A Rendszer Lényege: Egy teljesen új tervezésű, innovatív RF bajonett (lens mount) architektúra.",
                "Kompatibilitás: A meglévő EF és EF-S lencsék teljes körű támogatásának integrálása az új ökoszisztémába."
            ],
            "notes": "Üdvözlöm önöket! A mai prezentációban áttekintjük a Canon EOS R rendszerének stratégiai és műszaki alapjait. Látni fogjuk, hogy a jövőbeni piaci igények és az optikai határok feszegetése hogyan követelte meg egy teljesen új, kompromisszummentes mérnöki alap megteremtését."
        },
        {
            "title": "A Harmincéves EF Rendszer Korlátai",
            "bullets": [
                "Fizikai gátak: A hagyományos EF bajonett (54 mm átmérő, 44 mm távolság) elérte a lencsetervezési rugalmasság határait.",
                "Adatátviteli sebesség: A 8-tűs (pin) elektronikus kapcsolat sávszélessége mára már korlátozó tényezővé vált.",
                "Fejlesztési limitációk: Kevés az elektronikus csatorna az új, összetett operációs funkciók kiszolgálására.",
                "Fókusz korlátok: Szűkös lehetőségek a modern, szenzor-alapú autofókusz (AF) rendszerek maximalizálásában."
            ],
            "notes": "Bár a három évtizedes EF rendszer hatalmas nemzetközi siker volt, az olyan jövőbeni igények, mint a megnövelt szenzorfelbontás és az új videós funkciók, új architektúrát igényeltek. A régit nem lehetett tovább finomhangolni, radikális technológiai paradigmaváltásra volt szükségünk."
        },
        {
            "title": "Az RF Bajonett Innovációja",
            "bullets": [
                "Belső átmérő megőrzése: Az 54 mm-es belső átmérő megtartása az optimális fényáteresztés érdekében.",
                "Rövidített bázistávolság: A szenzor és a bajonett távolsága 44 mm-ről mindössze 20 mm-re csökkent.",
                "Növelt sávszélesség: A csatlakozók számának 8-ról 12-re történő emelése biztosítja a nagysebességű kommunikációt.",
                "Tervezési szabadság: Nagy átmérőjű hátsó lencsetagok helyezhetők közvetlenül a full-frame szenzor elé."
            ],
            "notes": "Ez a dia a rendszer szíve. Villamosmérnöki precizitással tervezték újra a kapcsolatot: a 12 aranyozott érintkező és az új protokollok olyan adatsűrűséget biztosítanak, amely valós idejű, masszív vezérlést tesz lehetővé a lencse és a váz között. A rövid, 20 mm-es távolság pedig az optikai tervezők számára eddig elérhetetlen szabadságot biztosít."
        },
        {
            "title": "Optikai Aberrációk Kezelése",
            "bullets": [
                "Fénytörési kihívások: A fénytörés fokozásával drasztikusan nőnek a kromatikus és monokromatikus képhibák (aberrációk).",
                "Optikai megoldás: A nagy hátsó lencsék \"kíméletesebb\" szögben vetítik a fénysugarakat a szenzor sarkaiba, csökkentve a torzítást.",
                "Digital Lens Optimizer (DLO): Beépített képoptimalizáló szoftverrendszer a kamerában a finomfelbontás javítására.",
                "DLO Működése: Valós időben, a RAW fájlok feldolgozásakor korrigálja a diffrakciót és a lencsespecifikus hibákat."
            ],
            "notes": "Minél jobban \"törjük\" a fényt, annál több a leküzdendő szférikus vagy asztigmatikus képhiba. Az új RF rendszer megengedi, hogy a lencsetagokkal enyhébb szögben vigyük be a fényt a szenzor széleire, ami radikálisan növeli a peremélességet. A DLO mindezt egy beépített, nagy teljesítményű szoftvermotorral egészíti ki."
        },
        {
            "title": "Az RF Lencsék Exkluzív Funkciói",
            "bullets": [
                "Vezérlőgyűrű (Control Ring): A fókuszgyűrű mellett egy új, testreszabható gyűrű az expozíciós beállításokhoz (Rekesz, ISO, Zársebesség).",
                "Változtatható fókusz-irány: A kézi élességállítás gyűrűjének forgásiránya szoftveresen, egyéni preferencia szerint megfordítható.",
                "\"Fly-by-wire\" fókuszálás: Nincs mechanikus kapcsolat; a fókuszgyűrű csak finom elektronikus jeleket küld az autófókusz motornak.",
                "Szenzorvédelem: Kikapcsoláskor a rekeszlamellák automatikusan bezárulnak, védve a szenzort a káros fénysugaraktól."
            ],
            "notes": "Az ergonómia új szintjét hozza el a programozható vezérlőgyűrű, amely \"kattanó\" visszajelzést ad forgatás közben. Mivel a fókuszálás mechanika helyett már teljesen \"fly-by-wire\" – elektronikus –, tökéletes, zökkenőmentes átmeneteket garantál, például professzionális videózás vagy Follow Focus rendszerek esetén."
        },
        {
            "title": "Az Induló RF Lencsekínálat (1. Rész)",
            "bullets": [
                "RF28-70mm F2 L USM: Szabványos, L-szériás zoomobjektív egyedülálló, állandó f/2.0 fényerővel a teljes tartományban.",
                "Méretbeli előny: A nagy bajonett és rövid bázistávolság miatt a frontlencse része kompaktabb maradhatott.",
                "RF50mm F1.2 L USM: Nagy átmérőjű, standard fix objektív extrém f/1.2 rekeszértékkel.",
                "Fejlett bevonatok: Air Sphere Coating (ASC) technológia a becsillanások és szellemkép minimálisra csökkentésére."
            ],
            "notes": "A 28-70-es f/2-es zoom megalkotása az EF rendszerben fizikai képtelenség lett volna ilyen méret- és súlyparaméterekkel. Az 50-es f/1.2 pedig a portréfotósok számára nyújt hihetetlen élességet már maximális nyitott rekesznél is, miközben rendkívül magas adatsebességgel kommunikál a vázzal."
        },
        {
            "title": "Az Induló RF Lencsekínálat (2. Rész)",
            "bullets": [
                "RF24-105mm F4 L IS USM: Kompakt f/4-es zoom, integrált képstabilizátorral professzionális és \"low-budget\" mozgóképes feladatokhoz.",
                "Nano USM meghajtás: Extra vékony, rendkívül gyors és néma ultrahangos fókuszmotor, ideális videózáshoz.",
                "RF35mm F1.8 MACRO IS STM: Széles látószögű objektív f/1.8 fényerővel és 0.5x makró nagyítással.",
                "Kompakt felépítés: Léptetőmotor (STM) gondoskodik a finom mozgásról, könnyű kivitelben."
            ],
            "notes": "Ez a két objektív a rendszer \"mindenese\". A 24-105-ös videósok kedvence lesz a szinte teljesen hangtalan és egyenletes \"start-and-stop\" fókuszálást biztosító Nano USM miatt. A 35mm-es pedig a kreatív \"street\" és makró fotózás tökéletes utazó eszköze, ami lenyűgöző 5-stop stabilizációt kínál."
        },
        {
            "title": "Rendszerkompatibilitás: EF Mount Adapterek",
            "bullets": [
                "Stratégiai integráció: A hatalmas globális EF/EF-S lencsekínálat maradéktalan átemelése az EOS R ökoszisztémába.",
                "Mount Adapter EF-EOS R: Alapmodell a tökéletes fizikai és elektronikus csatlakoztatáshoz, teljesítményveszteség nélkül.",
                "Control Ring Adapter: A régebbi EF lencséket is felruházza az új, beépített szenzoros expozíció-vezérlő gyűrű funkciójával.",
                "Drop-In Filter Adapter: Professzionális adapter beépített, cserélhető szűrő-nyílással (cirkuláris polár vagy változtatható ND)."
            ],
            "notes": "Természetesen nem hagyjuk cserben az évtizedek alatt felhalmozott EF lencse-arzenált. Épp ellenkezőleg: a speciális vezérlőgyűrűs és a bedobós szűrős adapterekkel a régi objektívek több funkciót kapnak az új vázakon, mint amivel eredetileg valaha is rendelkeztek."
        },
        {
            "title": "Képstabilizáció és Dual Pixel CMOS AF",
            "bullets": [
                "Interaktív Stabilizáció: A lencse giroszkópja és a szenzor elmozdulás-adatai (DIGIC 8 processzor) közösen harcolnak a remegés ellen.",
                "5-Tengelyes videós IS: A lencse optikai stabilizációját a szenzor elektronikus stabilizációja öttengelyessé egészíti ki mozgóképnél.",
                "Dual Pixel fókuszálás: Minden egyes képalkotó szenzorpixelben két fotodióda végzi a villámgyors fázisérzékelést.",
                "Kiterjesztett fókusz-lefedettség: A manuálisan választható fókuszpontok a kép szélességének 88%-át, magasságának 100%-át lefedik."
            ],
            "notes": "A kamera képfeldolgozó egysége szó szerint \"látja\" az elmosódást az új fókuszpixelek segítségével. A lencse és a váz valós időben dolgozzák fel ezeket az adatokat, így együttesen hoznak létre egy eddig sosem látott 5-tengelyes kompenzációs megoldást. A Dual Pixel fókusz is sokkal sűrűbb lett, kiváló találati aránnyal rossz fényviszonyok közt is."
        },
        {
            "title": "Összegzés: A Képalkotás Jövőjének Alapköve",
            "bullets": [
                "Mérnöki paradigmaváltás: Az RF rendszer az EF hagyatékára építve bontja le a lencsetervezés korábbi fizikai korlátait.",
                "Szinergia a rendszerben: Tökéletesebb hardveres integráció az optika, a mechanika és a nagysebességű digitális jelfeldolgozás között.",
                "Kész a jövő kihívásaira: Fejlesztési potenciál a folyamatosan növekvő szenzorfelbontáshoz és extenzív fényerő-igényekhez.",
                "Kompakt forma: Jelentős méretcsökkenés (flange back) anélkül, hogy a strukturális stabilitás vagy az optikai teljesítmény csorbát szenvedne."
            ],
            "notes": "Az EOS R megalkotásával a Canon nem csupán tartja a lépést a mirrorless trendekkel, hanem egy hosszú távú, robusztus technológiai platformot is nyújt a jövő számára. Kijelenthető, hogy ez a rendszer évtizedekig ki fogja szolgálni mind a technikai, mind a legmagasabb kreatív elvárásokat. Köszönöm megtisztelő figyelmüket!"
        }
    ]

    prs = Presentation()

    for slide_data in slides_data:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        # Cím beállítása
        title_shape = slide.shapes.title
        title_shape.text = slide_data["title"]

        # Tartalom doboz formázása
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        # Bekapcsoljuk az automatikus méretezést és a sortörést
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        for bullet in slide_data["bullets"]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            # Fixáljuk a betűméretet 20 pt-re, hogy biztosan elférjen
            p.font.size = Pt(20)

        # Előadói jegyzet
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = slide_data["notes"]

    filename = "Canon_EOS_R_Executive_Summary_v02.pptx"
    prs.save(filename)
    print(f"A prezentáció sikeresen elkészült és mentve lett: {filename}")

if __name__ == "__main__":
    create_presentation()